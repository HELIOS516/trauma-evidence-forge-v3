"""
Generate xABCDE Primary Survey PPTX from parsed slide content.
Design: Gamma Marine theme — dark navy bg, Calibri, accent #8C98CA
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ── Color palette ──────────────────────────────────────────────────────────
BG       = RGBColor(0x08, 0x0E, 0x26)   # dark navy
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BODY     = RGBColor(0xEB, 0xEC, 0xEF)
ACCENT   = RGBColor(0x8C, 0x98, 0xCA)   # light blue-purple
GRAY     = RGBColor(0x8C, 0x98, 0xCA)
HDR_BG   = RGBColor(0x1A, 0x23, 0x45)   # table header navy

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)


def make_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_layout(prs):
    """Return the blank slide layout."""
    return prs.slide_layouts[6]


def set_bg(slide, color=BG):
    """Fill slide background with solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height,
                text, font_name="Calibri", font_size=20,
                bold=False, italic=False, color=BODY,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_textbox_lines(slide, left, top, width, height,
                      lines, font_name="Calibri", font_size=20,
                      bold=False, italic=False, color=BODY,
                      align=PP_ALIGN.LEFT):
    """Add a textbox with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return txBox


def add_slide_number(slide, num):
    """Add slide number bottom-right."""
    add_textbox(slide, Inches(12.2), Inches(7.1), Inches(1.0), Inches(0.3),
                str(num), font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_bottom_line(slide, text):
    """Rounded rect at bottom with italic accent text."""
    left   = Inches(0.4)
    top    = Inches(6.55)
    width  = Inches(12.53)
    height = Inches(0.7)

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE — use int 1 for rounded rect via XML later
        left, top, width, height
    )
    # Set rounded corners via XML
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', 'val 25000')

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT

    line = shape.line
    line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_sources(slide, text):
    """Small gray sources line at very bottom."""
    add_textbox(slide, Inches(0.4), Inches(7.1), Inches(10.0), Inches(0.3),
                text, font_size=10, color=GRAY)


def add_key_stat(slide, text):
    """Large accent-colored key stat."""
    add_textbox(slide, Inches(0.4), Inches(5.3), Inches(12.5), Inches(0.9),
                text, font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════

def slide_title(prs, num):
    """Slide 1 — Title slide."""
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_textbox(slide,
        Inches(1.0), Inches(1.8), Inches(11.33), Inches(2.2),
        "xABCDE: Hemorrhage Kills First,\nSo Treat It First",
        font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        Inches(1.0), Inches(4.2), Inches(11.33), Inches(1.5),
        "Evan DeCan, MD  |  Division of Acute Care Surgery  |  University of Virginia  |  March 2026",
        font_size=18, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)

    add_slide_number(slide, num)
    return slide


def slide_learning_objectives(prs, num):
    """Slide 2 — Learning Objectives."""
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_textbox(slide, Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.7),
                "Learning Objectives", font_size=28, bold=True, color=WHITE)

    objectives = [
        "1.  Apply the xABCDE framework to systematically prioritize life-threatening interventions in a trauma patient",
        "2.  Identify clinical and physiologic indications for definitive airway management",
        "3.  Classify hemorrhagic shock by class and initiate massive transfusion protocol at appropriate thresholds",
    ]
    add_textbox_lines(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.5),
                      objectives, font_size=20, color=BODY)

    add_slide_number(slide, num)
    return slide


def slide_case(prs, num):
    """Slide 3 — Opening case."""
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_textbox(slide, Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.7),
                "22-Year-Old MVC: Unresponsive, Hypotensive, Tachycardic",
                font_size=26, bold=True, color=WHITE)

    bullets = [
        "Mechanism:   High-speed MVC, unrestrained driver",
        "Vitals:           HR 130  |  BP 80/50  |  RR 28  |  SpO2 92%",
        "GCS:              8  (E2 V2 M4)",
        "Key finding:  Distended abdomen, deformed pelvis",
        "Clinical Q:    What is your FIRST intervention?",
    ]
    add_textbox_lines(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(4.5),
                      bullets, font_size=20, color=BODY)

    add_bottom_line(slide, "Multiple life threats demand a systematic prioritized approach.")
    add_slide_number(slide, num)
    return slide


def content_slide(prs, num, title, bullets, key_stat, bottom_line, sources):
    """Generic content slide (Slides 4-7, 10-12, 14, 17-18)."""
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_textbox(slide, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.8),
                title, font_size=26, bold=True, color=WHITE)

    # Bullets
    bullet_top = Inches(1.15)
    bullet_h   = Inches(3.7)
    add_textbox_lines(slide, Inches(0.4), bullet_top, Inches(12.5), bullet_h,
                      ["•  " + b for b in bullets], font_size=19, color=BODY)

    # KEY STAT
    if key_stat:
        add_key_stat(slide, key_stat)

    # Bottom Line
    if bottom_line:
        add_bottom_line(slide, bottom_line)

    # Sources
    if sources:
        add_sources(slide, "Sources: " + sources)

    add_slide_number(slide, num)
    return slide


def slide_mcq(prs, num, title, question, options, bottom_line):
    """MCQ question slide."""
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_textbox(slide, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.6),
                title, font_size=26, bold=True, color=ACCENT)

    add_textbox(slide, Inches(0.4), Inches(1.0), Inches(12.5), Inches(1.5),
                question, font_size=20, bold=True, color=WHITE)

    opt_lines = ["  " + o for o in options]
    add_textbox_lines(slide, Inches(0.4), Inches(2.6), Inches(12.5), Inches(3.4),
                      opt_lines, font_size=20, color=BODY)

    add_bottom_line(slide, bottom_line)
    add_slide_number(slide, num)
    return slide


def slide_mcq_answer(prs, num, title, answer_bullets, bottom_line, sources=""):
    """MCQ answer slide."""
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_textbox(slide, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.8),
                title, font_size=24, bold=True, color=ACCENT)

    add_textbox_lines(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(4.8),
                      answer_bullets, font_size=19, color=BODY)

    add_bottom_line(slide, bottom_line)
    if sources:
        add_sources(slide, "Sources: " + sources)
    add_slide_number(slide, num)
    return slide


def slide_shock_table(prs, num):
    """Slide 13 — Hemorrhagic Shock Classification Table."""
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_textbox(slide, Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.65),
                "Class III and IV Hemorrhagic Shock Demand Immediate Blood Product Resuscitation",
                font_size=22, bold=True, color=WHITE)

    # Table data
    headers = ["Parameter", "Class I", "Class II", "Class III", "Class IV"]
    rows = [
        ["Blood loss",     "<750 mL\n(<15%)",       "750–1,500 mL\n(15–30%)",  "1,500–2,000 mL\n(30–40%)",  ">2,000 mL\n(>40%)"],
        ["Heart rate",     "<100",                  "100–120",                  "120–140",                    ">140"],
        ["Blood pressure", "Normal",                "Normal",                   "Decreased",                  "Decreased"],
        ["Resp rate",      "14–20",                 "20–30",                    "30–40",                      ">35"],
        ["Urine output",   ">30 mL/h",              "20–30 mL/h",              "5–15 mL/h",                  "Negligible"],
        ["GCS",            "15",                    "14–15",                    "12–14",                      "<12"],
        ["Treatment",      "Monitor",               "Possible transfusion",     "Blood products",             "Massive transfusion"],
    ]

    left = Inches(0.3)
    top  = Inches(1.0)
    w    = Inches(12.73)
    h    = Inches(5.7)

    table = slide.shapes.add_table(8, 5, left, top, w, h).table

    col_widths = [Inches(1.9), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.2)]
    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw

    def set_cell(cell, text, bg_color=None, font_color=WHITE,
                 font_size=13, bold=False, italic=False):
        cell.text = text
        tf = cell.text_frame
        tf.word_wrap = True
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.italic = italic
                run.font.color.rgb = font_color
        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color

    # Header row
    for col_idx, hdr in enumerate(headers):
        cell = table.cell(0, col_idx)
        set_cell(cell, hdr, bg_color=HDR_BG, font_color=WHITE, font_size=13, bold=True)

    # Data rows
    row_colors = [
        RGBColor(0x10, 0x18, 0x38),
        RGBColor(0x12, 0x1C, 0x3E),
        RGBColor(0x10, 0x18, 0x38),
        RGBColor(0x12, 0x1C, 0x3E),
        RGBColor(0x10, 0x18, 0x38),
        RGBColor(0x12, 0x1C, 0x3E),
        RGBColor(0x10, 0x18, 0x38),
    ]
    # Highlight Class III and IV columns (cols 3,4) in data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            highlight = col_idx in (3, 4)
            fc = ACCENT if highlight else BODY
            bc = row_colors[row_idx]
            set_cell(cell, cell_text, bg_color=bc, font_color=fc,
                     font_size=12, bold=(col_idx == 0))

    add_bottom_line(slide, "Our patient (HR 130, BP 80/50, GCS 8) fits Class III–IV — activate MTP.")
    add_sources(slide, "Sources: [18]")
    add_slide_number(slide, num)
    return slide


def slide_take_home(prs, num):
    """Slide 19 — Take-Home Pearls."""
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_textbox(slide, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.7),
                "Five Pearls That Will Save Lives on Your Next Trauma Activation",
                font_size=26, bold=True, color=WHITE)

    pearls = [
        "1.  X before A  —  hemorrhage control precedes airway in the exsanguinating patient",
        "2.  TXA within 3 hours  —  after 3 hours it increases mortality (CRASH-2)",
        "3.  Shock Index >1.0  —  calculate HR/SBP on every trauma patient; 4x mortality risk",
        "4.  Negative FAST is not reassuring  —  sensitivity only 74%; pursue further workup",
        "5.  Prevent the lethal triad  —  warm the patient, warm the blood, limit crystalloid",
    ]
    add_textbox_lines(slide, Inches(0.4), Inches(1.15), Inches(12.5), Inches(5.3),
                      pearls, font_size=20, color=BODY)

    add_bottom_line(slide, "xABCDE saves lives — treat first what kills first.")
    add_slide_number(slide, num)
    return slide


def slide_references(prs, num):
    """Slide 20 — References."""
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_textbox(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.55),
                "References", font_size=24, bold=True, color=WHITE)

    refs_col1 = [
        "1.  Netherton S, et al. Diagnostic accuracy of eFAST. CJEM. 2019;21(6):727-738. PMID: 31317856",
        "2.  Vang M, et al. Shock index as predictor of mortality. Eur J Trauma. 2022;48(4):2559-2566. PMID: 35258641",
        "3.  Shakur H, et al. CRASH-2. Lancet. 2010;376(9734):23-32. PMID: 20554319",
        "4.  Roberts I, et al. CRASH-2 exploratory. Lancet. 2011;377(9771):1096-1101. PMID: 21439633",
        "5.  Holcomb JB, et al. PROPPR. JAMA. 2015;313(5):471-482. PMID: 25647203",
        "6.  Sperry JL, et al. PAMPer. N Engl J Med. 2018;379(4):315-326. PMID: 30044935",
        "7.  Bickell WH, et al. Delayed fluid resuscitation. N Engl J Med. 1994;331(17):1105-1109. PMID: 7935634",
        "8.  Nunez TC, et al. ABC score. J Trauma. 2009;66(2):346-352. PMID: 19204506",
        "9.  Cotton BA, et al. Multicenter ABC validation. J Trauma. 2010;69(S1):S33-S39. PMID: 20622617",
        "10. Hatchimonji JS, et al. GCS=8 intubation dogma. Eur J Trauma. 2021;47:2073-2079. PMID: 32382780",
    ]
    refs_col2 = [
        "11. Eastridge BJ, et al. Death on battlefield 2001-2011. J Trauma. 2012;73(S5):S431-S437. PMID: 23192066",
        "12. Cannon JW, et al. DCR EAST guideline. J Trauma. 2017;82(3):605-617. PMID: 28225743",
        "13. LaGrone LN, et al. AAST/ACS-COT DCR protocol. J Trauma. 2024;96(3):510-520. PMID: 37697470",
        "14. Meizoso JP, et al. Whole blood EAST guideline. J Trauma. 2024;97(3):460-470. PMID: 38531812",
        "15. Mayglothling J, et al. Emergency intubation EAST. J Trauma. 2012;73(S4):S333-S340. PMID: 23114490",
        "16. Brown CVR, et al. WTA airway algorithm. Trauma Surg. 2020;5(1):e000539. PMID: 33083558",
        "17. de Moya M, et al. WTA pneumothorax algorithm. J Trauma. 2022;92(1):103-107. PMID: 34538823",
        "18. American College of Surgeons. ATLS 11th Edition. ACS; 2025.",
    ]

    add_textbox_lines(slide, Inches(0.3), Inches(0.9), Inches(6.3), Inches(6.5),
                      refs_col1, font_size=10, color=BODY)
    add_textbox_lines(slide, Inches(6.8), Inches(0.9), Inches(6.3), Inches(6.5),
                      refs_col2, font_size=10, color=BODY)

    add_slide_number(slide, num)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

def build_presentation():
    prs = make_prs()

    # ── Slide 1: Title ──────────────────────────────────────────────────
    slide_title(prs, 1)

    # ── Slide 2: Learning Objectives ────────────────────────────────────
    slide_learning_objectives(prs, 2)

    # ── Slide 3: Case ───────────────────────────────────────────────────
    slide_case(prs, 3)

    # ── Slide 4: X — Exsanguination ─────────────────────────────────────
    content_slide(prs, 4,
        title="Exsanguination Is the Leading Cause of Preventable Trauma Death",
        bullets=[
            "Hemorrhage causes 30–40% of trauma mortality",
            "90.9% of survivable battlefield deaths from bleeding",
            "Tourniquet within 15 min: high efficacy, low complication",
        ],
        key_stat="KEY STAT:  90.9% of preventable deaths = hemorrhage  (Eastridge 2012)",
        bottom_line="The X in xABCDE exists because hemorrhage kills fastest.",
        sources="[11][18]")

    # ── Slide 5: X — Hemorrhage Control Interventions ───────────────────
    content_slide(prs, 5,
        title="Massive External Hemorrhage Control Precedes Airway Management",
        bullets=[
            "Tourniquet for extremity hemorrhage",
            "Direct pressure and wound packing for junctional wounds",
            "Pelvic binder for unstable pelvic fractures",
        ],
        key_stat="KEY STAT:  Tourniquet + binder before intubation saves lives",
        bottom_line="In exsanguination, stop the bleeding before managing the airway.",
        sources="[11][18]")

    # ── Slide 6: A — Airway Algorithm ───────────────────────────────────
    content_slide(prs, 6,
        title="Airway Obstruction Remains a Top Preventable Prehospital Death After Hemorrhage",
        bullets=[
            "Jaw thrust, suction, remove debris first",
            "Supraglottic airway as bridge device",
            "RSI for definitive airway when indicated",
        ],
        key_stat="KEY STAT:  RSI is the recommended technique for trauma intubation  (EAST/WTA)",
        bottom_line="Assess, protect, then secure the airway using a stepwise algorithm.",
        sources="[15][16]")

    # ── Slide 7: A — GCS Nuance ─────────────────────────────────────────
    content_slide(prs, 7,
        title="GCS 8 or Less Is an Indication, Not an Absolute Mandate, for Intubation",
        bullets=[
            "GCS ≤8: strong indication for definitive airway",
            "Assess airway patency, gag reflex, trajectory",
            "Other triggers: obstruction, facial trauma, combativeness",
        ],
        key_stat="KEY STAT:  Rigid GCS ≤8 threshold may increase mortality  (PMID: 32382780)",
        bottom_line="GCS is one factor; clinical context determines the airway decision.",
        sources="[10][15][16]")

    # ── Slide 8: MCQ 1 ──────────────────────────────────────────────────
    slide_mcq(prs, 8,
        title="Clinical Decision Point",
        question="35M GSW right chest/abdomen. HR 140, BP 70/40, GCS 7. Active abdominal hemorrhage, gurgling respirations. First intervention?",
        options=[
            "A.  Rapid sequence intubation",
            "B.  Direct pressure to abdominal wound",
            "C.  Bilateral needle decompression",
            "D.  FAST exam",
        ],
        bottom_line="Think xABCDE — what kills this patient first?")

    # ── Slide 9: MCQ 1 Answer ───────────────────────────────────────────
    slide_mcq_answer(prs, 9,
        title="Answer: B — Control Hemorrhage First in the Exsanguinating Patient",
        answer_bullets=[
            "Why B:      Active exsanguination is the immediate lethal threat; X before A",
            "Why not A: GCS ≤8 alone does not override active hemorrhage control",
            "Why not C: No clinical signs of tension pneumothorax described",
            "Why not D: FAST delays definitive hemorrhage control in the unstable patient",
        ],
        bottom_line="Treat first what kills first — hemorrhage before airway in the exsanguinating patient.",
        sources="[11][18]")

    # ── Slide 10: B — Tension Pneumothorax ─────────────────────────────
    content_slide(prs, 10,
        title="Tension Pneumothorax Is a Clinical Diagnosis That Requires Immediate Decompression",
        bullets=[
            "Hypotension + absent breath sounds = decompress now",
            "Never delay for CXR or imaging",
            "Needle (2nd or 5th ICS) or finger thoracostomy",
        ],
        key_stat="KEY STAT:  Clinical diagnosis — do NOT wait for imaging",
        bottom_line="Suspected tension pneumothorax is treated, not imaged.",
        sources="[17]")

    # ── Slide 11: C — Blood on the Floor ───────────────────────────────
    content_slide(prs, 11,
        title='"Blood on the Floor and Four More" Localizes the Five Sources of Hemorrhage',
        bullets=[
            "Chest, abdomen, pelvis, long bones, external",
            "FAST identifies free fluid in unstable patients",
            "Pelvic X-ray and CXR guide intervention",
        ],
        key_stat="KEY STAT:  5 sources — chest, abdomen, pelvis, long bones, external",
        bottom_line="Identify the bleeding source to direct surgical or interventional control.",
        sources="[1][18]")

    # ── Slide 12: C — MTP Triggers ─────────────────────────────────────
    content_slide(prs, 12,
        title="Class III–IV Shock Requires Massive Transfusion Protocol Before Lab Results Return",
        bullets=[
            "Shock Index >1.0 = 4× mortality risk",
            "ABC score ≥2 activates MTP at bedside",
            "1:1:1 plasma : platelets : RBCs  (PROPPR)",
        ],
        key_stat="KEY STAT:  Shock Index >1.0 = 4-fold mortality increase  (Vang 2022)",
        bottom_line="Activate MTP on clinical criteria — do not wait for a hemoglobin.",
        sources="[2][5][8][9]")

    # ── Slide 13: C — Shock Classification Table ────────────────────────
    slide_shock_table(prs, 13)

    # ── Slide 14: C — FAST ──────────────────────────────────────────────
    content_slide(prs, 14,
        title="FAST Answers One Question: Is There Free Fluid in the Unstable Patient?",
        bullets=[
            "Sensitivity 74%, specificity 98% for free fluid",
            "Positive FAST + hypotension = operative intervention",
            "Negative FAST does NOT exclude injury",
        ],
        key_stat="KEY STAT:  Sensitivity only 74% — negative FAST does not rule out injury",
        bottom_line="FAST rules in; it does not rule out — a negative FAST demands further workup.",
        sources="[1]")

    # ── Slide 15: MCQ 2 ─────────────────────────────────────────────────
    slide_mcq(prs, 15,
        title="Management Decision",
        question="22M MVC patient: FAST positive (Morrison's pouch). BP 70/40 despite 2u pRBCs. GCS 7. Next step?",
        options=[
            "A.  CT abdomen/pelvis with contrast",
            "B.  Emergent exploratory laparotomy",
            "C.  Diagnostic peritoneal lavage",
            "D.  2 liters normal saline",
        ],
        bottom_line="FAST positive + unstable = OR, not CT scanner.")

    # ── Slide 16: MCQ 2 Answer ──────────────────────────────────────────
    slide_mcq_answer(prs, 16,
        title="Answer: B — FAST Positive and Hemodynamically Unstable Means Emergent Laparotomy",
        answer_bullets=[
            "Why B:      Positive FAST + refractory hypotension = OR for hemorrhage control",
            "Why not A: CT requires hemodynamic stability; this patient is in extremis",
            "Why not C: DPL is for equivocal FAST; this FAST is clearly positive",
            "Why not D: Crystalloid worsens the lethal triad in hemorrhagic shock",
        ],
        bottom_line="Positive FAST plus hemodynamic instability equals the operating room.",
        sources="[1][12][13]")

    # ── Slide 17: D — Disability ────────────────────────────────────────
    content_slide(prs, 17,
        title="GCS Components Plus Pupil Exam Constitute the Rapid Disability Assessment",
        bullets=[
            "GCS: Eye (1–4) + Verbal (1–5) + Motor (1–6) = 3–15",
            "Motor component best predictor of outcome",
            "Pupils: bilateral reactive vs fixed/dilated vs asymmetric",
        ],
        key_stat="KEY STAT:  GCS ≤8 = severe TBI; assess motor score specifically",
        bottom_line="Report GCS by components (E/V/M), not just the total — motor matters most.",
        sources="[10][18]")

    # ── Slide 18: E — Exposure ──────────────────────────────────────────
    content_slide(prs, 18,
        title="Exposure and Hypothermia Prevention Must Begin in the Trauma Bay",
        bullets=[
            "Full exposure: log roll for posterior exam",
            "Lethal triad: hypothermia + acidosis + coagulopathy",
            "Warm OR, warm blood, warm blankets",
        ],
        key_stat="KEY STAT:  47.8% mortality when lethal triad is complete",
        bottom_line="Undress the patient to find injuries, then actively prevent hypothermia.",
        sources="[12][13]")

    # ── Slide 19: Take-Home Pearls ──────────────────────────────────────
    slide_take_home(prs, 19)

    # ── Slide 20: References ─────────────────────────────────────────────
    slide_references(prs, 20)

    return prs


if __name__ == "__main__":
    out_path = "/Users/fee4jk/Documents/GitHub/own/trauma-evidence-forge-v3/projects/Pending/xabcde-primary-survey/xabcde-primary-survey-claude.pptx"
    prs = build_presentation()
    prs.save(out_path)
    import os
    size = os.path.getsize(out_path)
    print(f"Saved: {out_path}")
    print(f"File size: {size:,} bytes  ({size / 1024:.1f} KB)")
    print(f"Slides: {len(prs.slides)}")
